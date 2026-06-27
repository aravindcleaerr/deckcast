"""Export the rendered slide frames to other deck formats.

  pptx : one full-bleed slide per frame (16:9), narration as speaker notes.
  html : a single self-contained HTML deck (frames embedded), keyboard-navigable.

Both reuse the 1920x1080 PNG frames the pipeline already renders, so the exported
decks look exactly like the video — no separate layout engine.
"""
import base64
import html as _html
from pathlib import Path


def _existing(slides):
    return [s for s in slides if s.get("frame") and Path(s["frame"]).exists()]


def to_pptx(slides, dest, size, project="Deck"):
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        raise SystemExit("pptx export needs python-pptx — `pip install python-pptx` "
                         "(or `pip install \"deckcast[export]\"`).")
    slides = _existing(slides)
    if not slides:
        return None
    w, h = size
    prs = Presentation()
    # Match the frame aspect ratio on a standard widescreen canvas.
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(13.333 * h / w)
    blank = prs.slide_layouts[6]
    for s in slides:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(Path(s["frame"]).resolve()), 0, 0,
                                 width=prs.slide_width, height=prs.slide_height)
        note = s.get("narration")
        if note:
            slide.notes_slide.notes_text_frame.text = note
    Path(dest).parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dest))
    return dest


def to_html(slides, dest, size, project="Deck"):
    slides = _existing(slides)
    if not slides:
        return None
    secs = []
    for i, s in enumerate(slides):
        b64 = base64.b64encode(Path(s["frame"]).read_bytes()).decode()
        note = _html.escape(s.get("narration") or "")
        secs.append(
            f'<section class="slide" id="s{i+1}">'
            f'<img src="data:image/png;base64,{b64}" alt="Slide {i+1}"/>'
            f'<div class="notes">{note}</div></section>')
    out = (_HTML
           .replace("%%TITLE%%", _html.escape(project))
           .replace("%%TOTAL%%", str(len(slides)))
           .replace("%%SLIDES%%", "\n".join(secs)))
    Path(dest).parent.mkdir(parents=True, exist_ok=True)
    Path(dest).write_text(out, encoding="utf-8")
    return dest


_HTML = """<!DOCTYPE html>
<html lang="en"><head><meta charset="UTF-8"/>
<meta name="viewport" content="width=device-width, initial-scale=1"/>
<title>%%TITLE%%</title>
<style>
  *{margin:0;box-sizing:border-box}
  html,body{height:100%;background:#0c0f12;color:#e7ecf0;font-family:Inter,system-ui,sans-serif;overflow:hidden}
  #deck{height:100%;width:100%;display:flex;align-items:center;justify-content:center}
  .slide{display:none;width:100%;height:100%;align-items:center;justify-content:center}
  .slide.active{display:flex}
  .slide img{max-width:100%;max-height:100%;object-fit:contain;box-shadow:0 10px 60px rgba(0,0,0,.6)}
  .slide .notes{display:none}
  #bar{position:fixed;top:0;left:0;height:4px;background:#E8A33D;transition:width .2s;z-index:5}
  #hud{position:fixed;bottom:16px;right:22px;font-size:14px;color:#9aa7b2;z-index:5;
       letter-spacing:.04em;font-variant-numeric:tabular-nums}
  #nav{position:fixed;bottom:14px;left:22px;display:flex;gap:8px;z-index:5}
  #nav button{background:#1b2229;color:#cdd6df;border:1px solid #2b343d;border-radius:8px;
       padding:6px 12px;font-size:15px;cursor:pointer}
  #nav button:hover{background:#242d36}
  #notes{position:fixed;left:0;right:0;bottom:0;max-height:38vh;overflow:auto;
       background:rgba(8,11,14,.96);border-top:1px solid #2b343d;padding:20px 26px;
       font-size:18px;line-height:1.5;color:#dbe3ea;display:none;z-index:6}
  #notes.show{display:block}
  #notes b{color:#E8A33D;font-weight:600;display:block;margin-bottom:6px;font-size:13px;letter-spacing:.12em;text-transform:uppercase}
  #hint{position:fixed;top:14px;right:22px;font-size:12px;color:#5d6873;z-index:5}
</style></head>
<body>
  <div id="bar"></div>
  <div id="deck">%%SLIDES%%</div>
  <div id="hint">← → navigate · N notes · F fullscreen</div>
  <div id="nav"><button onclick="go(i-1)">‹ Prev</button><button onclick="go(i+1)">Next ›</button></div>
  <div id="hud"><span id="cur">1</span> / %%TOTAL%%</div>
  <div id="notes"><b>Speaker notes</b><span id="ntext"></span></div>
<script>
  var slides=[].slice.call(document.querySelectorAll('.slide'));
  var total=slides.length, i=0;
  function go(n){
    i=Math.max(0,Math.min(total-1,n));
    slides.forEach(function(s,k){s.classList.toggle('active',k===i)});
    document.getElementById('cur').textContent=i+1;
    document.getElementById('bar').style.width=((i+1)/total*100)+'%';
    var note=slides[i].querySelector('.notes');
    document.getElementById('ntext').textContent=note?note.textContent:'';
    if(history.replaceState) history.replaceState(null,'','#s'+(i+1));
  }
  function toggleNotes(){document.getElementById('notes').classList.toggle('show')}
  document.addEventListener('keydown',function(e){
    if(e.key==='ArrowRight'||e.key===' '||e.key==='PageDown'){go(i+1);e.preventDefault();}
    else if(e.key==='ArrowLeft'||e.key==='PageUp'){go(i-1);e.preventDefault();}
    else if(e.key==='Home'){go(0);} else if(e.key==='End'){go(total-1);}
    else if(e.key==='n'||e.key==='N'){toggleNotes();}
    else if(e.key==='f'||e.key==='F'){if(!document.fullscreenElement)document.documentElement.requestFullscreen();else document.exitFullscreen();}
  });
  document.getElementById('deck').addEventListener('click',function(e){
    go(e.clientX < window.innerWidth/2 ? i-1 : i+1);
  });
  var h=parseInt((location.hash.match(/\\d+/)||[1])[0],10)-1;
  go(isNaN(h)?0:h);
</script>
</body></html>"""
